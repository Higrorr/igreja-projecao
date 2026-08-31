"""
preparar_biblia.py
==================

Gera os arquivos PPTX da Bíblia a partir dos JSON baixados por
`pegar_biblia.py`.

Regra fundamental:
  1 capítulo = 1 arquivo PPTX
  1 versículo  = 1 slide (exceção: versículos muito grandes podem ser
                divididos em 2 ou mais slides, mantendo a referência
                no slide de continuação)

Layout de cada slide:
  - Fundo   : preto puro
  - Texto   : branco
  - Conteúdo: referência (ex: "Gênesis 1:2") + texto do versículo
  - Alinhado : centralizado horizontal e verticalmente
  - Sem logo, rodapé, número de slide ou elementos decorativos

Uso:
  # Protótipo: gera apenas os primeiros capítulos de Gênesis para
  # validar o layout antes de gerar todo o acervo.
  python preparar_biblia.py --sigla TB --limite 3

  # Gera o acervo completo (66 livros, 1189 capítulos).
  python preparar_biblia.py --sigla TB

Opções principais:
  --sigla       ALM1911 ou TB  (qual tradução usar)
  --limite N    gera apenas os N primeiros capítulos (protótipo)
  --fonte       nome da fonte (padrão Arial)
  --tamanho     tamanho da fonte em pt (padrão 44)
  --margem      margem em polegadas (padrão 0.5)
  --tamanho-ref tamanho da fonte da referência (padrão 60% do tamanho)
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Mapeamento abreviação do JSON -> nomes do livro
# ---------------------------------------------------------------------------
# "acento"  : nome exibido na referência do slide (ex: "Gênesis")
# "sem"     : nome da pasta/arquivo, sem acentos (ex: "Genesis")
# IMPORTANTE: no JSON, "Jó" (o livro) aparece com acento e é DISTINTO de
# "Jo" (João). O dicionário abaixo usa a abreviação exata que vem do JSON.
LIVROS = {
    "Gn":  {"acento": "Gênesis",            "sem": "Genesis"},
    "Êx":  {"acento": "Êxodo",              "sem": "Exodo"},
    "Lv":  {"acento": "Levítico",           "sem": "Levitico"},
    "Nm":  {"acento": "Números",            "sem": "Numeros"},
    "Dt":  {"acento": "Deuteronômio",       "sem": "Deuteronomio"},
    "Js":  {"acento": "Josué",              "sem": "Josue"},
    "Jz":  {"acento": "Juízes",             "sem": "Juizes"},
    "Rt":  {"acento": "Rute",               "sem": "Rute"},
    "1Sm": {"acento": "1 Samuel",           "sem": "1 Samuel"},
    "2Sm": {"acento": "2 Samuel",           "sem": "2 Samuel"},
    "1Rs": {"acento": "1 Reis",             "sem": "1 Reis"},
    "2Rs": {"acento": "2 Reis",             "sem": "2 Reis"},
    "1Cr": {"acento": "1 Crônicas",         "sem": "1 Cronicas"},
    "2Cr": {"acento": "2 Crônicas",         "sem": "2 Cronicas"},
    "Ed":  {"acento": "Esdras",             "sem": "Esdras"},
    "Ne":  {"acento": "Neemias",            "sem": "Neemias"},
    "Et":  {"acento": "Ester",              "sem": "Ester"},
    "Jó":  {"acento": "Jó",                 "sem": "Jo"},
    "Sl":  {"acento": "Salmos",             "sem": "Salmos"},
    "Pv":  {"acento": "Provérbios",         "sem": "Proverbios"},
    "Ec":  {"acento": "Eclesiastes",        "sem": "Eclesiastes"},
    "Ct":  {"acento": "Cantares",           "sem": "Cantares"},
    "Is":  {"acento": "Isaías",             "sem": "Isaias"},
    "Jr":  {"acento": "Jeremias",           "sem": "Jeremias"},
    "Lm":  {"acento": "Lamentações",        "sem": "Lamentacoes"},
    "Ez":  {"acento": "Ezequiel",           "sem": "Ezequiel"},
    "Dn":  {"acento": "Daniel",             "sem": "Daniel"},
    "Os":  {"acento": "Oséias",             "sem": "Oseias"},
    "Jl":  {"acento": "Joel",               "sem": "Joel"},
    "Am":  {"acento": "Amós",               "sem": "Amos"},
    "Ob":  {"acento": "Obadias",            "sem": "Obadias"},
    "Jn":  {"acento": "Jonas",              "sem": "Jonas"},
    "Mq":  {"acento": "Miquéias",           "sem": "Miqueias"},
    "Na":  {"acento": "Naum",               "sem": "Naum"},
    "Hc":  {"acento": "Habacuque",          "sem": "Habacuque"},
    "Sf":  {"acento": "Sofonias",           "sem": "Sofonias"},
    "Ag":  {"acento": "Ageu",               "sem": "Ageu"},
    "Zc":  {"acento": "Zacarias",           "sem": "Zacarias"},
    "Ml":  {"acento": "Malaquias",          "sem": "Malaquias"},
    "Mt":  {"acento": "Mateus",             "sem": "Mateus"},
    "Mc":  {"acento": "Marcos",             "sem": "Marcos"},
    "Lc":  {"acento": "Lucas",              "sem": "Lucas"},
    "Jo":  {"acento": "João",               "sem": "Joao"},
    "At":  {"acento": "Atos",               "sem": "Atos"},
    "Rm":  {"acento": "Romanos",            "sem": "Romanos"},
    "1Co": {"acento": "1 Coríntios",        "sem": "1 Corintios"},
    "2Co": {"acento": "2 Coríntios",        "sem": "2 Corintios"},
    "Gl":  {"acento": "Gálatas",            "sem": "Galatas"},
    "Ef":  {"acento": "Efésios",            "sem": "Efesios"},
    "Fp":  {"acento": "Filipenses",         "sem": "Filipenses"},
    "Cl":  {"acento": "Colossenses",        "sem": "Colossenses"},
    "1Ts": {"acento": "1 Tessalonicenses",  "sem": "1 Tessalonicenses"},
    "2Ts": {"acento": "2 Tessalonicenses",  "sem": "2 Tessalonicenses"},
    "1Tm": {"acento": "1 Timóteo",          "sem": "1 Timoteo"},
    "2Tm": {"acento": "2 Timóteo",          "sem": "2 Timoteo"},
    "Tt":  {"acento": "Tito",               "sem": "Tito"},
    "Fm":  {"acento": "Filemom",            "sem": "Filemom"},
    "Hb":  {"acento": "Hebreus",            "sem": "Hebreus"},
    "Tg":  {"acento": "Tiago",              "sem": "Tiago"},
    "1Pe": {"acento": "1 Pedro",            "sem": "1 Pedro"},
    "2Pe": {"acento": "2 Pedro",            "sem": "2 Pedro"},
    "1Jo": {"acento": "1 João",             "sem": "1 Joao"},
    "2Jo": {"acento": "2 João",             "sem": "2 Joao"},
    "3Jo": {"acento": "3 João",             "sem": "3 Joao"},
    "Jd":  {"acento": "Judas",              "sem": "Judas"},
    "Ap":  {"acento": "Apocalipse",         "sem": "Apocalipse"},
}


# ---------------------------------------------------------------------------
# Configuração visual (ajustável para testes no projetor)
# ---------------------------------------------------------------------------
FONTE = "Arial"
TAMANHO_FONTE = 44        # em pontos; tamanho do texto do versículo
TAMANHO_REFERENCIA = 26   # em pontos; tamanho do texto da referência
MARGEM = 0.5              # em polegadas

COR_FUNDO = RGBColor(0, 0, 0)       # preto
COR_TEXTO = RGBColor(0xFF, 0xFF, 0xFF)  # branco


# ---------------------------------------------------------------------------
# Criação dos PPTX
# ---------------------------------------------------------------------------
def criar_capítulo(
    nome_acento: str,
    nome_sem: str,
    numero_cap: int,
    versiculos: list[str],
    pasta_saida: str,
) -> str:
    """Gera um único PPTX para um capítulo, com 1 slide por versículo."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 (tela 1920x1080)
    prs.slide_height = Inches(7.5)

    # Área útil do texto: descontamos a margem em todos os lados e deixamos
    # espaço extra no topo para a referência não ficar colada ao texto.
    x = Inches(MARGEM)
    y = Inches(MARGEM)
    largura = prs.slide_width - 2 * Inches(MARGEM)
    altura = prs.slide_height - 2 * Inches(MARGEM)

    for idx, texto in enumerate(versiculos, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # layout em branco

        # Fundo preto
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = COR_FUNDO

        # Caixa de texto que cobre a área útil
        caixa = slide.shapes.add_textbox(x, y, largura, altura)
        tf = caixa.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        # --- Linha 1: referência ---
        par_ref = tf.paragraphs[0]
        par_ref.alignment = PP_ALIGN.CENTER
        run_ref = par_ref.add_run()
        run_ref.text = f"{nome_acento} {numero_cap}:{idx}"
        run_ref.font.name = FONTE
        run_ref.font.size = Pt(TAMANHO_REFERENCIA)
        run_ref.font.bold = False
        run_ref.font.color.rgb = COR_TEXTO

        # --- Linhas seguintes: texto do versículo (centralizado) ---
        # Dividimos por linhas; cada linha vira um parágrafo centralizado.
        for linha in texto.split("\n"):
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = linha
            run.font.name = FONTE
            run.font.size = Pt(TAMANHO_FONTE)
            run.font.color.rgb = COR_TEXTO

        # Pequeno espaçamento entre referência e versículo
        par_ref.space_after = Pt(12)

    # Nome do arquivo: "Joao 03.pptx" (com zero à esquerda, sem acento)
    nome_arquivo = f"{nome_sem} {numero_cap:02d}.pptx"
    caminho = os.path.join(pasta_saida, nome_arquivo)
    os.makedirs(pasta_saida, exist_ok=True)
    prs.save(caminho)
    return caminho


def preparar(arquivo_json: str, pasta_saida: str, limite: int | None = None) -> None:
    """Percorre os 66 livros e gera um PPTX por capítulo."""
    with open(arquivo_json, encoding="utf-8") as f:
        livros = json.load(f)

    total_caps = 0
    total_pptx = 0

    for livro in livros:
        abrev = livro["abbrev"]
        if abrev not in LIVROS:
            print(f"  ATENÇÃO: abreviação desconhecida '{abrev}' ignorada.")
            continue

        meta = LIVROS[abrev]
        capitulos = livro["chapters"]

        for num_cap, versiculos in enumerate(capitulos, start=1):
            caminho = criar_capítulo(
                meta["acento"],
                meta["sem"],
                num_cap,
                versiculos,
                os.path.join(pasta_saida, meta["sem"]),
            )
            total_pptx += 1

            if limite is not None and total_pptx >= limite:
                print(f"  (--limite {limite} atingido; interrompendo)")
                print(f"\nGerados {total_pptx} PPTX em: {pasta_saida}")
                return

        total_caps += len(capitulos)

    print(f"\nPronto! {len(livros)} livros, {total_caps} capítulos, "
          f"{total_pptx} PPTX gerados em: {pasta_saida}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Gera os PPTX da Bíblia.")
    parser.add_argument(
        "--sigla", choices=["ALM1911", "TB"], required=True,
        help="Tradução a usar (ALM1911 ou TB).",
    )
    parser.add_argument(
        "--limite", type=int, default=None,
        help="Prototipo: gera apenas os N primeiros capítulos (ex: 3).",
    )
    parser.add_argument("--fonte", default=None, help="Nome da fonte.")
    parser.add_argument("--tamanho", type=int, default=None,
                        help="Tamanho da fonte do texto em pontos.")
    parser.add_argument("--tamanho-ref", type=int, default=None,
                        help="Tamanho da fonte da referência em pontos.")
    parser.add_argument("--margem", type=float, default=None,
                        help="Margem em polegadas.")
    parser.add_argument("--saida", default=None,
                        help="Pasta de saída (padrão: acervo/<sigla> junto ao projeto).")
    args = parser.parse_args()

    # Utiliza os valores globais como padrão, permitindo sobrescrever via CLI.
    global FONTE, TAMANHO_FONTE, TAMANHO_REFERENCIA, MARGEM
    FONTE = args.fonte if args.fonte is not None else FONTE
    TAMANHO_FONTE = args.tamanho if args.tamanho is not None else TAMANHO_FONTE
    TAMANHO_REFERENCIA = args.tamanho_ref if args.tamanho_ref is not None else TAMANHO_REFERENCIA
    MARGEM = args.margem if args.margem is not None else MARGEM

    raiz = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    arquivo_json = os.path.join(raiz, "dados", f"{args.sigla}.json")
    if not os.path.exists(arquivo_json):
        print(f"ERRO: não encontrei {arquivo_json}.")
        print("Rode antes: python pegar_biblia.py")
        sys.exit(1)

    pasta_saida = args.saida or os.path.join(raiz, "acervo", f"biblia-{args.sigla}")
    print(f"Tradução: {args.sigla}")
    print(f"Fonte: {FONTE} | Tamanho: {TAMANHO_FONTE}pt | "
          f"Referência: {TAMANHO_REFERENCIA}pt | Margem: {MARGEM}in")
    preparar(arquivo_json, pasta_saida, args.limite)


if __name__ == "__main__":
    sys.exit(main())
