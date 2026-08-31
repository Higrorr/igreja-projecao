"""
preparar_harpa.py
=================

Gera os arquivos PPTX padronizados dos 640 hinos da Harpa Cristã, a partir
dos arquivos .ppt originais baixados pelo usuário.

Os arquivos .ppt originais são formato binário legado que o python-pptx NÃO
consegue ler. Para extrair o texto, este script usa a interface COM do
PowerPoint (pywin32), que está instalado neste computador (o PC do projetor).

Regras de layout (iguais às da Bíblia):
  - Slide 16:9, fundo preto, texto branco, centralizado, margem configurável
  - Sem logo, rodapé, número de slide ou elementos decorativos

Estrutura de cada hino (escolhas do usuário):
  - Slide 1 : título do hino  (ex: "Harpa 133" + nome do hino, em destaque)
  - Slide 2+: cada estrofe em um slide, com referência pequena "Harpa <n>"
              no topo (igual à Bíblia exibe a referência do versículo)
  - Refrão : mantido exatamente como no original (sem duplicação automática)

Uso:
  # Protótipo: gera apenas os primeiros N hinos para validar o layout.
  python preparar_harpa.py --origem <pasta> --limite 3

  # Completo: gera os 640 hinos.
  python preparar_harpa.py --origem <pasta>

Onde <pasta> é a pasta com os arquivos .ppt (ex: harpa/.../HC640 Para DataShow).
A saída vai para acervo/harpa/01xxx.pptx por padrão.
"""

import argparse
import glob
import os
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Configuração visual (mesmos valores iniciais da Bíblia)
# ---------------------------------------------------------------------------
FONTE = "Arial"
TAMANHO_TITULO = 60        # tamanho do slide de título
TAMANHO_FONTE = 44         # tamanho do texto da estrofe
TAMANHO_REFERENCIA = 26    # tamanho da referência "Harpa <n>"
MARGEM = 0.5               # em polegadas

COR_FUNDO = RGBColor(0, 0, 0)
COR_TEXTO = RGBColor(0xFF, 0xFF, 0xFF)


def extrair_hino(caminho_ppt, app):
    """
    Abre um .ppt via COM (somente leitura) e devolve:
        (titulo, [estrofe1, estrofe2, ...])
    O título vem do slide 1; as estrofes dos slides 2..N.

    Caso o hino tenha apenas 1 slide, o título e a letra estão juntos nele.
    Nesse caso separamos o título (antes da dupla quebra de linha) da letra,
    e dividimos a letra em estrofes por blocos de linhas.
    """
    pres = app.Presentations.Open(caminho_ppt, ReadOnly=1, WithWindow=0)
    try:
        def texto_slide(i):
            partes = []
            for sh in pres.Slides(i).Shapes:
                if sh.HasTextFrame:
                    t = sh.TextFrame.TextRange.Text or ""
                    # \x0b é a quebra de linha usada dentro de uma caixa de
                    # texto no PowerPoint. Normalizamos para \n.
                    partes.append(t.replace("\x0b", "\n").strip())
            return "\n".join(partes)

        n_slides = pres.Slides.Count

        if n_slides >= 2:
            titulo = texto_slide(1)
            estrofes = []
            for i in range(2, n_slides + 1):
                estrofes.append(texto_slide(i))
            return titulo, estrofes

        # Caso de 1 slide: título + letra juntos.
        texto = texto_slide(1)
        # Separa título (antes da dupla quebra de linha) do conteúdo.
        parte_titulo = texto
        parte_letra = ""
        if "\n\n" in texto:
            parte_titulo, parte_letra = texto.split("\n\n", 1)
        elif not texto.strip():
            return "", []
        # Quando não há dupla quebra, assume que a primeira linha é o título.
        else:
            linhas = [l for l in texto.split("\n") if l.strip()]
            if linhas:
                parte_titulo = linhas[0]
                parte_letra = "\n".join(linhas[1:])

        # Divide a letra em estrofes separadas por linhas em branco.
        estrofes = [b.strip() for b in parte_letra.split("\n\n") if b.strip()]
        return parte_titulo, estrofes
    finally:
        pres.Close()


def montar_titulo(titulo_raw):
    """
    O slide de título costuma ser algo como "133\x0bNO ROL DO LIVRO".
    Devolve (nome_do_hino, tem_numero). O número já vem embutido,
    então apenas separamos para ter o texto do título puro.
    """
    import unicodedata
    # Remove espaços não-separáveis e normaliza espaços.
    def limpar(s):
        s = s.replace("\xa0", " ").replace("\u200b", "")
        return " ".join(s.split())

    linhas = [limpar(l) for l in titulo_raw.split("\n") if limpar(l)]
    if not linhas:
        return "Harpa", ""
    # Tenta extrair um número da primeira linha
    m = re.match(r"^(\d+)\s*[:\-–]?\s*(.*)$", linhas[0])
    if m:
        num = m.group(1)
        resto = m.group(2).strip()
        nome = resto if resto else " ".join(linhas[1:]).strip()
        return nome, num
    # Sem número na primeira linha: usa tudo
    return " ".join(linhas), ""


def adicionar_slide(prs, texto_principal, tamanho_principal, ref_linha):
    """Cria um slide 16:9 preto com a referência (se houver) e o texto centralizado."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COR_FUNDO

    x = Inches(MARGEM)
    y = Inches(MARGEM)
    largura = prs.slide_width - 2 * Inches(MARGEM)
    altura = prs.slide_height - 2 * Inches(MARGEM)

    caixa = slide.shapes.add_textbox(x, y, largura, altura)
    tf = caixa.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    if ref_linha:
        par_ref = tf.paragraphs[0]
        par_ref.alignment = PP_ALIGN.CENTER
        run = par_ref.add_run()
        run.text = ref_linha
        run.font.name = FONTE
        run.font.size = Pt(TAMANHO_REFERENCIA)
        run.font.color.rgb = COR_TEXTO
        par_ref.space_after = Pt(12)
        primeiro = 1
    else:
        primeiro = 0

    # Texto principal: cada linha visual vira um parágrafo centralizado.
    linhas = texto_principal.split("\n")
    for idx, linha in enumerate(linhas):
        if idx == 0 and primeiro == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = linha
        run.font.name = FONTE
        run.font.size = Pt(tamanho_principal)
        run.font.color.rgb = COR_TEXTO


def nome_arquivo_seguro(numero, nome_hino):
    """
    Monta o nome do arquivo com o padrão "<numero:03d> - <nome>.pptx".
    Remove caracteres inválidos para nomes de arquivo no Windows e
    normaliza espaços/tamanho para evitar nomes excessivamente longos.
    """
    import unicodedata

    def limpar(texto):
        # Espaços não-separáveis -> espaço comum
        texto = texto.replace("\xa0", " ")
        # Remove caracteres inválidos no Windows
        for c in '<>:"/\\|?*':
            texto = texto.replace(c, "")
        texto = " ".join(texto.split())
        return texto.strip()

    nome = limpar(nome_hino)
    if len(nome) > 60:  # limita o nome para o arquivo não ficar enorme
        nome = nome[:60].rstrip()
    base = f"{numero:03d} - {nome}" if nome else f"{numero:03d}"
    return f"{base}.pptx"


def gerar_hino(numero, titulo, estrofes, pasta_saida):
    """Gera um único .pptx para o hino, com título + 1 estrofe por slide."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: título do hino
    nome, _ = titulo
    slide_titulo = f"Harpa {numero}"
    adicionar_slide(prs, slide_titulo + ("\n" + nome if nome else ""),
                    TAMANHO_TITULO, ref_linha="")

    # Slides seguintes: cada estrofe, com referência "Harpa <n>" no topo
    for estrofe in estrofes:
        if not estrofe:
            continue
        adicionar_slide(prs, estrofe, TAMANHO_FONTE, ref_linha=f"Harpa {numero}")

    os.makedirs(pasta_saida, exist_ok=True)
    nome_arquivo = nome_arquivo_seguro(numero, nome)
    caminho = os.path.join(pasta_saida, nome_arquivo)
    prs.save(caminho)
    return caminho


def main():
    parser = argparse.ArgumentParser(description="Gera os PPTX padronizados da Harpa.")
    parser.add_argument("--origem", default=None,
                        help="Pasta com os arquivos .ppt originais. Se omitida, "
                             "procurará automaticamente dentro de harpa/.")
    parser.add_argument("--limite", type=int, default=None,
                        help="Prototipo: processa apenas os N primeiros hinos.")
    parser.add_argument("--fonte", default=None, help="Nome da fonte (padrão Arial).")
    parser.add_argument("--tamanho", type=int, default=None,
                        help="Tamanho do texto da estrofe em pontos.")
    parser.add_argument("--tamanho-titulo", type=int, default=None,
                        help="Tamanho do título em pontos.")
    parser.add_argument("--tamanho-ref", type=int, default=None,
                        help="Tamanho da referência em pontos.")
    parser.add_argument("--margem", type=float, default=None,
                        help="Margem em polegadas.")
    parser.add_argument("--saida", default=None,
                        help="Pasta de saída (padrão: acervo/harpa junto ao projeto).")
    args = parser.parse_args()

    global FONTE, TAMANHO_FONTE, TAMANHO_TITULO, TAMANHO_REFERENCIA, MARGEM
    if args.fonte is not None:
        FONTE = args.fonte
    if args.tamanho is not None:
        TAMANHO_FONTE = args.tamanho
    if args.tamanho_titulo is not None:
        TAMANHO_TITULO = args.tamanho_titulo
    if args.tamanho_ref is not None:
        TAMANHO_REFERENCIA = args.tamanho_ref
    if args.margem is not None:
        MARGEM = args.margem

    if not args.origem:
        raiz = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        pasta_harpa = os.path.join(raiz, "harpa")
        if os.path.isdir(pasta_harpa):
            for root, _dirs, files in os.walk(pasta_harpa):
                if any(f.lower().endswith(".ppt") for f in files):
                    args.origem = root
                    break
    if not args.origem or not os.path.isdir(args.origem):
        print("ERRO: não encontrei pasta de origem com arquivos .ppt.")
        print("Passe --origem <caminho> ou coloque os .ppt dentro de harpa/ do projeto.")
        sys.exit(1)

    arquivos = sorted(glob.glob(os.path.join(args.origem, "*.ppt")))
    if not arquivos:
        print("ERRO: nenhum arquivo .ppt encontrado na pasta.")
        sys.exit(1)

    # Ordena numericamente pelo número inicial do nome do arquivo.
    def num_arquivo(p):
        m = re.match(r"(\d+)", os.path.basename(p))
        return int(m.group(1)) if m else 10**9
    arquivos.sort(key=num_arquivo)

    if args.limite:
        arquivos = arquivos[:args.limite]

    raiz = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    pasta_saida = args.saida or os.path.join(raiz, "acervo", "harpa")

    import win32com.client
    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = 1
    try:
        for caminho in arquivos:
            numero = num_arquivo(caminho)
            titulo, estrofes = extrair_hino(caminho, app)
            titulo_info = montar_titulo(titulo)
            destino = gerar_hino(numero, titulo_info, estrofes, pasta_saida)
            print(f"  {numero:3d} - {titulo_info[0][:40]:40} -> {os.path.basename(destino)} "
                  f"({len(estrofes)} estrofes)")
    finally:
        try:
            app.Quit()
        except Exception:
            # A instância do PowerPoint pode já ter sido encerrada; ignora.
            pass

    print(f"\nPronto! {len(arquivos)} hinos gerados em: {pasta_saida}")


if __name__ == "__main__":
    sys.exit(main())
